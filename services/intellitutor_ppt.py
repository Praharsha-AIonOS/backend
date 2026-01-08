# PPT parsing and summarization for IntelliTutor (Feature 4)
import os
import subprocess
import fitz
from pptx import Presentation
import requests
from dotenv import load_dotenv

load_dotenv()

GROQ_ENDPOINT = "https://api.groq.com/openai/v1/chat/completions"


def _summarize_with_groq(text: str) -> str:
    api_key = os.getenv("GROQ_API_KEY")
    if not api_key or not text.strip():
        print("🔴 Groq API key missing or text is empty")
        return text

    payload = {
        "model": "llama-3.1-8b-instant",
        "messages": [
            {
                "role": "system",
                "content": (
                    "Produce a clear, concise summary of the provided content. "
                    "Write it in a natural speaking pace with short, smooth sentences. "
                    "Do NOT include phrases like 'here is the summary', "
                    "'explanation', or 'narration'. "
                    "Only output the summary text."
                )
            },
            {"role": "user", "content": text}
        ],
        "max_tokens": 200,
        "temperature": 0.3,
    }
    try:
        r = requests.post(
            GROQ_ENDPOINT,
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
            },
            json=payload,
            timeout=60,
        )

        if r.status_code != 200:
            return text

        return r.json()["choices"][0]["message"]["content"].strip()
    except Exception as e:
        print(f"🔴 Groq API error: {e}")
        return text


def parse_ppt(ppt_path: str, output_dir: str):
    """
    1. Extract text from PPT slides
    2. Summarize each slide via Groq
    3. Convert PPT -> PDF using LibreOffice
    4. Render each slide as PNG
    """
    os.makedirs(output_dir, exist_ok=True)

    # -----------------------------
    # Extract raw slide text
    # -----------------------------
    prs = Presentation(ppt_path)
    raw_texts = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        raw_texts.append("\n".join(parts))

    # -----------------------------
    # Summarize slide text
    # -----------------------------
    slide_texts = []
    for idx, text in enumerate(raw_texts):
        logger.info("Summarizing slide %d", idx + 1)
        slide_texts.append(_summarize_with_groq(text))

    # -----------------------------
    # Convert PPT -> PDF via libreoffice
    # -----------------------------
    subprocess.run(
        ["soffice", "--headless", "--convert-to", "pdf", ppt_path, "--outdir", output_dir],
        check=True,
        timeout=120,
    )

    pdf_name = os.path.splitext(os.path.basename(ppt_path))[0] + ".pdf"
    pdf_path = os.path.join(output_dir, pdf_name)

    if not os.path.exists(pdf_path):
        raise RuntimeError("PDF conversion failed")

    # -----------------------------
    # Render slides as images
    # -----------------------------
    doc = fitz.open(pdf_path)

    slide_images = []
    for i in range(len(doc)):
        page = doc.load_page(i)
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_path = os.path.join(output_dir, f"slide_{i}.png")
        pix.save(img_path)
        slide_images.append(img_path)

    return {
        "slide_images": slide_images,
        "slide_texts": slide_texts,
    }

